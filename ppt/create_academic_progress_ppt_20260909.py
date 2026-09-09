#!/usr/bin/env python3
"""Build the editable 2026-09-09 RL-GenRisk academic progress deck.

All scientific plots are rendered as high-resolution PNG files and embedded
in the PPTX.  Titles, explanations, annotations, and diagram elements remain
editable in PowerPoint/WPS.  The deck uses Train/Validation artifacts only;
Test labels are not read by this script.
"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, "/tmp/rlgenrisk_pptx")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches

from create_project_progress_ppt import (
    BLUE,
    GRID,
    H,
    INK,
    LIGHT,
    MUTED,
    NAVY,
    ORANGE,
    PALE_BLUE,
    PALE_RED,
    PALE_TEAL,
    RED,
    TEAL,
    W,
    WHITE,
    add_rich_text,
    add_text,
    bullet,
    image,
    line,
    metric_card,
    rect,
    rgb,
)


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = Path(__file__).resolve().parent
ASSET_DIR = PPT_DIR / "assets_20260909_v2_通俗讲解版"
OUT = PPT_DIR / "RL_GenRisk_学术汇报_20260909_v2_通俗讲解版.pptx"

FONT = "Microsoft YaHei"
COLORS = {
    "navy": "#102A43",
    "blue": "#1479B8",
    "teal": "#008B8B",
    "orange": "#E08A31",
    "red": "#C94C4C",
    "muted": "#627D98",
    "grid": "#D9E2EC",
    "light": "#F4F7FA",
}


def chart_style():
    chinese_font = Path("/mnt/c/Windows/Fonts/msyh.ttc")
    if chinese_font.exists():
        font_manager.fontManager.addfont(str(chinese_font))
        chart_font = font_manager.FontProperties(fname=str(chinese_font)).get_name()
    else:
        chart_font = "DejaVu Sans"
    plt.rcParams.update(
        {
            "font.family": chart_font,
            "font.size": 9,
            "axes.unicode_minus": False,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.labelcolor": COLORS["navy"],
            "axes.titlecolor": COLORS["navy"],
            "xtick.color": COLORS["muted"],
            "ytick.color": COLORS["muted"],
            "axes.edgecolor": COLORS["grid"],
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.facecolor": "white",
        }
    )


def save_fig(fig, name: str) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    fig.savefig(path, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def ordered_preferences(frame: pd.DataFrame) -> pd.DataFrame:
    if "w_discovery" not in frame.columns:
        frame = frame.copy()
        frame["w_discovery"] = frame.Preference.str.extract(r"_d([0-9.]+)")[0].astype(float)
    return frame.sort_values("w_discovery").reset_index(drop=True)


def build_charts() -> dict[str, Path]:
    chart_style()
    charts: dict[str, Path] = {}

    # 1. Mutation-frequency label bias. These values are from the completed
    # audit recorded in README; Test is intentionally omitted from the plot.
    fig, ax = plt.subplots(figsize=(7.1, 3.15))
    labels = ["全部基因背景\n中位数", "训练标签\n(n=16)", "验证标签\n(n=25)"]
    vals = [50.0, 99.8, 98.5]
    bars = ax.bar(labels, vals, color=["#B8C8D8", COLORS["blue"], COLORS["teal"]], width=0.58)
    ax.set_ylim(0, 108)
    ax.set_ylabel("在 9,039 个基因中的突变频率百分位")
    ax.set_title("已知驱动基因标签集中在高突变频率端", loc="left", weight="bold")
    ax.grid(axis="y", color="#E6EDF3", linewidth=0.8)
    ax.set_axisbelow(True)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, val + 2.2, f"{val:.1f}%", ha="center", weight="bold", color=COLORS["navy"])
    charts["label_bias"] = save_fig(fig, "01_label_frequency_bias.png")

    # 2. RL necessity: DDQN and independent contextual bandit.
    rn = pd.read_csv(ROOT / "outputs/rlnecessity_20260907_152513_comparison/rl_necessity_summary.csv")
    rn = rn[rn.Method.isin(["ddqn", "contextual_bandit"])].copy()
    rn["w_discovery"] = rn.Preference.str.extract(r"_d([0-9.]+)")[0].astype(float)
    fig, axes = plt.subplots(1, 2, figsize=(8.0, 3.3), sharex=True)
    for method, color, label in [
        ("ddqn", COLORS["orange"], "双深度Q网络（DDQN）"),
        ("contextual_bandit", COLORS["blue"], "即时价值模型"),
    ]:
        part = ordered_preferences(rn[rn.Method.eq(method)])
        x = part.w_discovery.to_numpy()
        axes[0].errorbar(x, part["NDCG@150_mean"], yerr=part["NDCG@150_sd"], marker="o", lw=2.1, capsize=3, color=color, label=label)
        axes[1].errorbar(x, part["DiscoveryFoldEnrichment@150_mean"], yerr=part["DiscoveryFoldEnrichment@150_sd"], marker="o", lw=2.1, capsize=3, color=color, label=label)
    axes[0].set_title("找回已知基因", loc="left", weight="bold")
    axes[0].set_ylabel("NDCG@150（均值 ± 标准差）")
    axes[1].set_title("发现低频新候选", loc="left", weight="bold")
    axes[1].set_ylabel("富集倍数（均值 ± 标准差）")
    for ax in axes:
        ax.set_xlabel("新候选发现目标的权重")
        ax.set_xticks([0, .2, .5, .8, 1])
        ax.grid(color="#E6EDF3", linewidth=.8)
        ax.set_axisbelow(True)
    axes[0].legend(frameon=False, fontsize=8)
    fig.suptitle("去掉“下一状态估值”后，性能没有一致下降", x=.07, ha="left", weight="bold", color=COLORS["navy"])
    fig.tight_layout()
    charts["rl_necessity"] = save_fig(fig, "02_ddqn_vs_bandit.png")

    # 3. History ablation heat maps.
    hist = pd.read_csv(ROOT / "outputs/historyablation_20260907_200151_comparison/history_ablation_summary.csv")
    modes = ["full", "no_history", "shuffled_history"]
    prefs = ["r1.00_d0.00", "r0.80_d0.20", "r0.50_d0.50", "r0.20_d0.80", "r0.00_d1.00"]
    labels = ["0.0", "0.2", "0.5", "0.8", "1.0"]
    matrices = []
    for metric in ["NDCG@150_mean", "DiscoveryFoldEnrichment@150_mean"]:
        p = hist.pivot(index="HistoryMode", columns="Preference", values=metric).reindex(index=modes, columns=prefs)
        matrices.append(p.to_numpy())
    fig, axes = plt.subplots(1, 2, figsize=(8.0, 3.2))
    titles = ["已知基因找回：NDCG@150", "新候选发现：富集倍数"]
    cmaps = ["Blues", "GnBu"]
    for ax, mat, title, cmap in zip(axes, matrices, titles, cmaps):
        im = ax.imshow(mat, aspect="auto", cmap=cmap)
        ax.set_title(title, loc="left", weight="bold")
        ax.set_xticks(range(5), labels)
        ax.set_xlabel("新候选发现目标的权重")
        ax.set_yticks(range(3), ["完整历史", "无历史", "打乱历史"])
        for i in range(mat.shape[0]):
            for j in range(mat.shape[1]):
                ax.text(j, i, f"{mat[i,j]:.2f}", ha="center", va="center", fontsize=8, color="#102A43")
        fig.colorbar(im, ax=ax, fraction=.045, pad=.03)
    fig.suptitle("真实的已选基因历史没有稳定优于两种消融", x=.07, ha="left", weight="bold", color=COLORS["navy"])
    fig.tight_layout()
    charts["history"] = save_fig(fig, "03_history_ablation.png")

    # 4. MORL frontier coverage across architectural repairs.
    sources = [
        ("单值Q\n共享模型", ROOT / "outputs/morl_shared_50ep_20260903_165939_comparison/morl_vs_scalar_frontier_coverage.csv", COLORS["blue"]),
        ("双目标Q", ROOT / "outputs/morl_vectorq_50ep_20260904_160342_comparison/morl_vs_scalar_frontier_coverage.csv", COLORS["orange"]),
        ("双目标Q\n+ 尺度校准", ROOT / "outputs/morl_vectorq_popart_50ep_20260904_201812_comparison/morl_vs_scalar_frontier_coverage.csv", COLORS["red"]),
    ]
    counts = []
    for label, path, color in sources:
        f = pd.read_csv(path)
        counts.append((label, int(f.morl_covers_scalar.sum()), color))
    fig, ax = plt.subplots(figsize=(7.0, 3.2))
    bars = ax.bar([x[0] for x in counts], [x[1] for x in counts], color=[x[2] for x in counts], width=.55)
    ax.set_ylim(0, 15.5)
    ax.set_yticks([0, 5, 10, 15])
    ax.set_ylabel("覆盖的独立权重基准点（共15个）")
    ax.set_title("多次修复后，共享多目标模型的覆盖能力仍未稳定", loc="left", weight="bold")
    ax.grid(axis="y", color="#E6EDF3")
    ax.set_axisbelow(True)
    for b, (_, value, _) in zip(bars, counts):
        ax.text(b.get_x()+b.get_width()/2, value+.5, f"{value}/15", ha="center", weight="bold")
    charts["morl_coverage"] = save_fig(fig, "04_morl_frontier_coverage.png")

    # 5. Five-seed shared preference-conditioned scalar bandit.
    shared = pd.read_csv(ROOT / "outputs/preference_bandit_20260908_190650_5seed_extension/preference_bandit_summary.csv")
    shared = ordered_preferences(shared)
    fig, axes = plt.subplots(1, 2, figsize=(8.0, 3.25), sharex=True)
    x = shared.w_discovery.to_numpy()
    axes[0].errorbar(x, shared["NDCG@150_mean"], yerr=shared["NDCG@150_sd"], marker="o", lw=2.4, capsize=3, color=COLORS["blue"])
    axes[0].errorbar(x, shared["Recall@150_mean"], yerr=shared["Recall@150_sd"], marker="s", lw=2.0, capsize=3, color="#60A5C2", label="Recall")
    axes[0].set_title("已知基因找回能力下降", loc="left", weight="bold")
    axes[0].set_ylabel("指标值")
    axes[0].legend(["排序质量（NDCG）", "找回比例（Recall）"], frameon=False, fontsize=8)
    axes[1].errorbar(x, shared["DiscoveryPrecision@150_mean"], yerr=shared["DiscoveryPrecision@150_sd"], marker="o", lw=2.4, capsize=3, color=COLORS["teal"])
    axes[1].set_title("新候选证据精度上升", loc="left", weight="bold")
    axes[1].set_ylabel("发现精度")
    for ax in axes:
        ax.set_xlabel("新候选发现目标的权重")
        ax.set_xticks([0, .2, .5, .8, 1])
        ax.grid(color="#E6EDF3", linewidth=.8)
        ax.set_axisbelow(True)
    fig.suptitle("单输出共享即时模型：5个随机种子均复现总体权衡方向", x=.07, ha="left", weight="bold", color=COLORS["navy"])
    fig.tight_layout()
    charts["shared_tradeoff"] = save_fig(fig, "05_shared_bandit_tradeoff.png")

    # 6. Dual-head seen-preference objective curves.
    dual = pd.read_csv(ROOT / "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_summary.csv")
    dual_seen = ordered_preferences(dual[dual.Scope.eq("seen")])
    fig, axes = plt.subplots(2, 2, figsize=(8.0, 4.8), sharex=True)
    specs = [
        ("排序质量（NDCG@150）", "NDCG@150_mean", "NDCG@150_sd", COLORS["blue"]),
        ("找回比例（Recall@150）", "Recall@150_mean", "Recall@150_sd", "#60A5C2"),
        ("新候选发现精度", "DiscoveryPrecision@150_mean", "DiscoveryPrecision@150_sd", COLORS["teal"]),
        ("新候选富集倍数", "DiscoveryFoldEnrichment@150_mean", "DiscoveryFoldEnrichment@150_sd", COLORS["orange"]),
    ]
    x = dual_seen.w_discovery.to_numpy()
    for ax, (title, mean, sd, color) in zip(axes.flat, specs):
        ax.errorbar(x, dual_seen[mean], yerr=dual_seen[sd], marker="o", capsize=3, lw=2.2, color=color)
        ax.set_title(title, loc="left", weight="bold", fontsize=10)
        ax.grid(color="#E6EDF3", linewidth=.8)
        ax.set_axisbelow(True)
    for ax in axes[-1]:
        ax.set_xlabel("新候选发现目标的权重")
        ax.set_xticks([0, .2, .5, .8, 1])
    fig.suptitle("双输出即时模型恢复了清晰的两目标权衡（3个随机种子）", x=.07, ha="left", weight="bold", color=COLORS["navy"])
    fig.tight_layout()
    charts["dual_summary"] = save_fig(fig, "06_dual_head_summary.png")

    # 7. Dual-head trajectories at seed level, including unseen preferences.
    per = pd.read_csv(ROOT / "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_per_seed.csv")
    fig, ax = plt.subplots(figsize=(7.3, 3.5))
    for seed, group in per.groupby("Seed"):
        group = group.sort_values("w_discovery")
        ax.plot(group["NDCG@150"], group["DiscoveryPrecision@150"], marker="o", lw=1.9, label=f"随机种子 {seed}")
        for _, row in group.iterrows():
            if row.w_discovery in (0.0, .5, 1.0) or np.isclose(row.w_discovery, .9):
                ax.annotate(f"{row.w_discovery:.1f}", (row["NDCG@150"], row["DiscoveryPrecision@150"]), xytext=(3,3), textcoords="offset points", fontsize=7)
    ax.set_xlabel("已知基因找回：NDCG@150")
    ax.set_ylabel("新候选发现精度")
    ax.set_title("偏好权重会改变排序，但纯发现端对随机种子敏感", loc="left", weight="bold")
    ax.grid(color="#E6EDF3", linewidth=.8)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8)
    charts["dual_seed"] = save_fig(fig, "07_dual_head_seed_trajectories.png")

    # 8. Smoothness comparison and bracket success rates.
    smooth = pd.read_csv(ROOT / "outputs/dual_head_bandit_20260908_210039_comparison/interpolation_smoothness_comparison.csv")
    methods = ["单输出\n共享模型", "双输出模型"]
    smooth = smooth.set_index("Method").loc[["single_head_shared_bandit", "dual_head_bandit"]]
    metrics = [
        ("相邻权重\n名单重合度", "AdjacentTop150JaccardMean"),
        ("NDCG落在\n相邻区间", "NDCG@150_within_bracket_rate"),
        ("Recall落在\n相邻区间", "Recall@150_within_bracket_rate"),
        ("发现精度落在\n相邻区间", "DiscoveryPrecision@150_within_bracket_rate"),
    ]
    x = np.arange(len(metrics)); width=.34
    fig, ax = plt.subplots(figsize=(7.5, 3.3))
    a = [smooth.iloc[0][c] for _, c in metrics]
    b = [smooth.iloc[1][c] for _, c in metrics]
    ax.bar(x-width/2, a, width, label=methods[0], color="#8EBAD4")
    ax.bar(x+width/2, b, width, label=methods[1], color=COLORS["orange"])
    ax.set_xticks(x, [m[0] for m in metrics])
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("比例 / 相似度")
    ax.set_title("双输出增强了权衡，但没有改善中间偏好的平滑性", loc="left", weight="bold")
    ax.grid(axis="y", color="#E6EDF3")
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8)
    for i, vals in enumerate(zip(a, b)):
        for offset, val in zip([-width/2, width/2], vals):
            ax.text(i+offset, val+.025, f"{val:.2f}", ha="center", fontsize=8)
    charts["smoothness"] = save_fig(fig, "08_interpolation_smoothness.png")

    # 9. Adjacent-policy overlap heat map exposes seed48 endpoint jump.
    jac = pd.read_csv(ROOT / "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_adjacent_jaccard.csv")
    jac["interval"] = jac.LeftPreference.str.replace("r", "", regex=False).str.replace("_d", "/", regex=False) + "->" + jac.RightPreference.str.replace("r", "", regex=False).str.replace("_d", "/", regex=False)
    intervals = jac[jac.Seed.eq(42)].interval.tolist()
    mat = jac.pivot(index="Seed", columns="interval", values="Top150Jaccard").reindex(index=[42,45,48], columns=intervals).to_numpy()
    fig, ax = plt.subplots(figsize=(8.2, 2.9))
    im = ax.imshow(mat, aspect="auto", cmap="RdYlGn", vmin=.2, vmax=1.0)
    ax.set_yticks(range(3), ["随机种子42", "随机种子45", "随机种子48"])
    short = ["0-.1", ".1-.2", ".2-.35", ".35-.5", ".5-.65", ".65-.8", ".8-.9", ".9-1"]
    ax.set_xticks(range(8), short)
    ax.set_xlabel("相邻的新候选发现权重区间")
    ax.set_title("随机种子48在纯发现端附近出现Top-150名单突变", loc="left", weight="bold")
    for i in range(mat.shape[0]):
        for j in range(mat.shape[1]):
            ax.text(j, i, f"{mat[i,j]:.2f}", ha="center", va="center", fontsize=8)
    fig.colorbar(im, ax=ax, fraction=.025, pad=.02, label="名单重合度")
    charts["jaccard_heatmap"] = save_fig(fig, "09_dual_head_adjacent_jaccard.png")

    # 10. Paired dual-head minus independent scalar bandit comparison (n=2).
    paired = pd.read_csv(ROOT / "outputs/dual_head_bandit_20260908_210039_comparison/dual_minus_scalar_bandit_paired.csv")
    agg = paired.groupby("Preference", as_index=False).mean(numeric_only=True)
    pref_order = ["r1.00_d0.00", "r0.80_d0.20", "r0.50_d0.50", "r0.20_d0.80", "r0.00_d1.00"]
    agg = agg.set_index("Preference").loc[pref_order].reset_index()
    x = np.arange(5)
    fig, axes = plt.subplots(1, 2, figsize=(8.0, 3.15))
    vals1 = agg["dual_minus_scalar_bandit_NDCG@150"]
    vals2 = agg["dual_minus_scalar_bandit_DiscoveryPrecision@150"]
    for ax, vals, title, color in [
        (axes[0], vals1, "已知基因NDCG差值", COLORS["blue"]),
        (axes[1], vals2, "新候选发现精度差值", COLORS["teal"]),
    ]:
        ax.axhline(0, color=COLORS["muted"], lw=.9)
        ax.bar(x, vals, color=[color if v >= 0 else COLORS["red"] for v in vals])
        ax.set_xticks(x, ["0", ".2", ".5", ".8", "1"])
        ax.set_xlabel("新候选发现目标的权重")
        ax.set_title(title, loc="left", weight="bold")
        ax.grid(axis="y", color="#E6EDF3")
        ax.set_axisbelow(True)
    fig.suptitle("双输出减去独立权重即时模型：权衡更强，但并非全面更好（配对n=2）", x=.07, ha="left", weight="bold", color=COLORS["navy"])
    fig.tight_layout()
    charts["paired_scalar"] = save_fig(fig, "10_dual_minus_scalar.png")

    # 11. Three fixed operating profiles proposed from seen preferences.
    profiles = dual_seen[dual_seen.Preference.isin(["r0.80_d0.20", "r0.50_d0.50", "r0.20_d0.80"])].copy()
    profiles = profiles.set_index("Preference").loc[["r0.80_d0.20", "r0.50_d0.50", "r0.20_d0.80"]]
    fig, ax = plt.subplots(figsize=(6.8, 3.4))
    ax.scatter(profiles["NDCG@150_mean"], profiles["DiscoveryFoldEnrichment@150_mean"], s=[125,155,185], c=[COLORS["blue"], COLORS["teal"], COLORS["orange"]], edgecolor="white", linewidth=1.5)
    for label, (_, row) in zip(["已知基因优先", "均衡", "新候选优先"], profiles.iterrows()):
        ax.annotate(label, (row["NDCG@150_mean"], row["DiscoveryFoldEnrichment@150_mean"]), xytext=(7,7), textcoords="offset points", fontsize=9, weight="bold")
    ax.axhline(1.0, color=COLORS["red"], ls="--", lw=1, label="没有富集")
    ax.set_xlabel("已知基因找回：NDCG@150")
    ax.set_ylabel("新候选发现：富集倍数")
    ax.set_title("当前证据更支持三档固定策略，而不是连续调节", loc="left", weight="bold")
    ax.grid(color="#E6EDF3")
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper right")
    charts["profiles"] = save_fig(fig, "11_three_fixed_profiles.png")

    return charts


def add_text_local(slide, text, x, y, w, h, size=18, color=NAVY, bold=False,
                   align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    return add_text(slide, text, x, y, w, h, size, color, bold, align, valign, name=FONT)


def base_slide(prs, section: str, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    section_names = {
        "Executive summary": "阶段结论",
        "Background": "研究背景",
        "Why two objectives": "为什么需要双目标",
        "Initial hypothesis": "最初假设",
        "Evidence-led route": "证据驱动的路线",
        "Problem 1": "问题一",
        "Problem 2": "问题二",
        "Problem 3": "问题三",
        "Current method": "当前方法",
        "Positive result 1": "阶段成果一",
        "Positive result 2": "阶段成果二",
        "Robustness": "稳定性检查",
        "Failure boundary": "尚未解决的问题",
        "Model comparison": "模型对比",
        "Current decision": "当前决策",
        "Expected outcomes": "预期成果",
        "Next work": "后续工作",
        "Claims & limitations": "结论边界",
        "Selected references": "参考文献",
    }
    add_text_local(slide, section_names.get(section, section), .56, .28, 3.2, .22, 8.5, BLUE, True)
    line(slide, .56, .58, 12.76, .58, GRID, .75)
    add_text_local(slide, title, .56, .73, 12.0, .48, 24, NAVY, True)
    if subtitle:
        add_text_local(slide, subtitle, .58, 1.23, 12.0, .28, 10.5, MUTED)
    line(slide, .56, 7.10, 12.76, 7.10, GRID, .55)
    add_text_local(slide, "RL-GenRisk  |  仅训练集/验证集  |  测试集标签冻结", .56, 7.18, 8.0, .16, 7.5, MUTED)
    add_text_local(slide, f"{len(prs.slides):02d}", 12.28, 7.15, .45, .18, 8, MUTED, True, align=PP_ALIGN.RIGHT)
    return slide


def source_link(slide, label: str, target: str, x=.68, y=6.72, w=11.2):
    box = add_text_local(slide, f"数据源：{label}", x, y, w, .18, 8.2, MUTED)
    run = box.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = "file:///E:/Projects/RL-GenRisk-main/" + target.replace("\\", "/")
    return box


def external_link(slide, label: str, url: str, x, y, w, h=.26, size=11.5):
    box = add_text_local(slide, label, x, y, w, h, size, BLUE, False)
    run = box.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = url
    return box


def simplify_language(prs):
    """Replace repeated technical English in editable slide text.

    Hyperlinked source paths and paper titles are left unchanged so the links
    remain recognizable and copyable.
    """
    replacements = [
        ("Recovery–Discovery", "已知基因找回–新候选发现"),
        ("Recovery-Discovery", "已知基因找回–新候选发现"),
        ("Recovery-heavy", "已知基因优先"),
        ("Discovery-heavy", "新候选优先"),
        ("contextual bandit", "即时价值模型"),
        ("scalar bandit", "独立权重即时模型"),
        ("shared bandit", "共享即时模型"),
        ("bandit", "即时价值模型"),
        ("scalar Pareto frontier", "独立权重模型的帕累托前沿"),
        ("scalar frontier", "独立权重模型前沿"),
        ("vector-Q / PopArt", "双目标Q与尺度校准"),
        ("vector-Q", "双目标Q"),
        ("bootstrap", "下一状态估值"),
        ("Full history", "完整历史"),
        ("No-history", "无历史"),
        ("Shuffled-history", "打乱历史"),
        ("true history", "真实历史"),
        ("real history", "真实历史"),
        ("history", "已选基因历史"),
        ("preference conditioning", "偏好权重控制"),
        ("preference", "偏好权重"),
        ("trade-off", "此消彼长的权衡"),
        ("greedy rollout", "逐步贪心生成排序"),
        ("one-pass", "一次性打分"),
        ("magnitude dominance", "数值尺度支配"),
        ("Fold enrichment", "富集倍数"),
        ("enrichment", "富集程度"),
        ("Precision", "精度"),
        ("Fold", "富集倍数"),
        ("gradient cosine", "梯度方向相似度"),
        ("conditioning", "偏好条件输入"),
        ("collapse", "条件失效"),
        ("shared trunk", "共享主干网络"),
        ("trunk", "主干网络"),
        ("双 head", "双输出"),
        ("head", "输出头"),
        ("Top-150 Jaccard", "Top-150名单重合度"),
        ("Jaccard", "名单重合度"),
        ("scalarized", "独立权重"),
        ("scalarization", "加权合成"),
        ("scalar Q", "单值Q"),
        ("scalar 权重", "单独权重"),
        ("scalar 前沿", "独立权重前沿"),
        ("Q-learning", "Q值学习"),
        ("rollout", "逐步生成排序"),
        ("cosine", "余弦相似度"),
        ("TD", "目标误差"),
        ("Balanced", "均衡"),
        ("context", "上下文信息"),
        ("序列 RL", "序列强化学习"),
        ("RL 回路", "强化学习回路"),
        ("scalar", "独立权重"),
        ("static mutation", "静态突变频率排序"),
        ("mutation", "突变频率"),
        ("baseline", "基线"),
        ("driver", "驱动基因"),
        ("reward", "训练奖励"),
        ("ranking", "排序"),
        ("Spearman", "单调趋势系数"),
        ("w_disc", "发现权重"),
        ("Train", "训练集"),
        ("shared", "共享"),
        ("Recovery", "已知基因找回"),
        ("Discovery", "新候选发现"),
        ("Validation-only", "仅使用验证集"),
        ("Validation", "验证集"),
        ("Test", "测试集"),
        ("seed", "随机种子"),
        ("endpoint", "端点"),
        ("paired", "配对"),
    ]
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink.address:
                        continue
                    value = run.text
                    for old, new in replacements:
                        value = value.replace(old, new)
                    run.text = value


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(WHITE)
    rect(s, 0, 0, .24, H, BLUE, BLUE, radius=False)
    add_text_local(s, "RL-GenRisk", .74, .72, 3.0, .34, 14, BLUE, True)
    add_text_local(s, "从序列强化学习到\n可按目标偏好切换的基因排序", .74, 1.32, 9.6, 1.34, 34, NAVY, True)
    add_text_local(s, "阶段性学术汇报：路线依据、问题诊断、实验证据与预期成果", .77, 2.98, 9.7, .34, 15, INK)
    line(s, .77, 3.62, 11.92, 3.62, GRID, .8)
    add_text_local(s, "核心结论", .77, 3.92, 1.3, .22, 11, TEAL, True)
    add_text_local(s, "两个目标之间确实存在取舍；序列强化学习的必要性不足；\n当前证据支持三档固定策略，不支持连续调节目标偏好。", .77, 4.32, 8.8, .78, 20, NAVY, True)
    rect(s, 10.13, 4.20, 2.16, 1.52, PALE_TEAL, None)
    add_text_local(s, "仅使用验证集", 10.37, 4.56, 1.68, .25, 11, TEAL, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "测试集冻结", 10.37, 5.02, 1.68, .25, 16, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "2026-09-09", .77, 6.53, 2.2, .20, 10, MUTED)


def add_slides(prs, charts):
    # 2. Plain-language metric guide.
    s = base_slide(prs, "先看懂指标", "四个核心指标分别在回答什么？", "后续所有图都只是在这四个问题上比较模型。")
    cards = [
        ("排序质量\nNDCG@150", "不仅看是否命中，还看命中的已知基因排得是否靠前。\n越高越好。", BLUE),
        ("找回比例\nRecall@150", "Top-150 找回了多少验证集已知基因。\n验证集共25个，每多命中1个就增加0.04。", "3FA8C9"),
        ("新候选精度\n高证据比例", "低频、未标注候选中，有多少具备预先定义的高证据。\n越高越好。", TEAL),
        ("富集倍数\n相对随机水平", "相对于候选池随机水平，高证据候选增加了多少倍。\n大于1才表示富集。", ORANGE),
    ]
    xs = [.66, 3.73, 6.80, 9.87]
    for (head, note, color), x in zip(cards, xs):
        rect(s, x, 1.77, 2.54, 3.80, LIGHT, None)
        rect(s, x, 1.77, 2.54, .12, color, color, radius=False)
        add_text_local(s, head, x+.19, 2.16, 2.16, .63, 17, color, True, align=PP_ALIGN.CENTER)
        line(s, x+.28, 3.04, x+2.26, 3.04, GRID, .8)
        add_text_local(s, note, x+.25, 3.37, 2.04, 1.48, 13.2, INK, align=PP_ALIGN.CENTER)
    rect(s, .68, 5.91, 11.72, .51, PALE_RED, None)
    add_text_local(s, "注意：这些都是验证集指标。它们用于比较方案，不等同于临床有效性或实验室证实。", .94, 6.06, 11.20, .22, 12.5, RED, True, align=PP_ALIGN.CENTER)

    # 3. Plain-language model guide.
    s = base_slide(prs, "先看懂模型", "几个反复出现的模型术语，用一句话怎么理解？")
    terms = [
        ("双深度Q网络（DDQN）", "不仅学习当前动作的奖励，还估计它对后续选择的影响。", "包含“下一状态估值”", BLUE),
        ("即时价值模型", "根据当前信息直接判断某个基因值不值得选。", "不估计未来回报", TEAL),
        ("多目标强化学习（MORL）", "同时处理“找回已知”和“发现新候选”两个互相冲突的目标。", "理想上输出多种取舍", ORANGE),
        ("偏好权重 w", "两个数字表示使用者更重视哪个目标，例如(0.8,0.2)。", "数字之和为1", NAVY),
    ]
    for i, (head, meaning, key, color) in enumerate(terms):
        y = 1.66 + i*1.17
        rect(s, .72, y, 3.17, .86, LIGHT, None)
        rect(s, .72, y, .08, .86, color, color, radius=False)
        add_text_local(s, head, .98, y+.16, 2.63, .48, 14.2, color, True, align=PP_ALIGN.CENTER)
        add_text_local(s, meaning, 4.16, y+.11, 5.68, .51, 13.5, INK)
        rect(s, 10.17, y+.10, 2.14, .55, PALE_BLUE if color != ORANGE else "FFF4E5", None)
        add_text_local(s, key, 10.32, y+.24, 1.84, .22, 11.2, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "随机种子：改变模型初始化和采样随机性的编号；多个随机种子方向一致，结果才更可信。", .74, 6.40, 11.60, .25, 12, MUTED, align=PP_ALIGN.CENTER)

    # 2. Executive summary.
    s = base_slide(prs, "Executive summary", "一句话状态：路线已经从“修强化学习”转为“验证更简单的排序模型”")
    metric_card(s, .62, 1.69, 2.92, 1.65, "已确认", "权衡存在", "Recovery 权重下降时，Discovery 指标整体上升。", TEAL)
    metric_card(s, 3.73, 1.69, 2.92, 1.65, "已否定", "序列优势不足", "DDQN 未稳定优于 bandit；真实 history 也未稳定胜出。", ORANGE, warning=True)
    metric_card(s, 6.84, 1.69, 2.92, 1.65, "最新结果", "双 head 有效", "Recovery/Discovery 分离更清楚，但插值未更平滑。", BLUE)
    metric_card(s, 9.95, 1.69, 2.74, 1.65, "当前决策", "固定三档", "不扩 5 seed；不继续追求连续 preference。", RED, warning=True)
    add_text_local(s, "阶段性贡献", .67, 3.92, 2.0, .25, 13, NAVY, True)
    bullet(s, "建立了可审计的 Train/Validation/Test 隔离与统一 greedy rollout 评价协议。", .69, 4.36, 11.6, .40, 15)
    bullet(s, "通过消融逐步排除了 reward 调参、bootstrap 和历史状态作为主要收益来源。", .69, 4.95, 11.6, .40, 15)
    bullet(s, "得到一个可解释的双即时价值模型，并明确识别其连续控制与跨 seed 稳定性边界。", .69, 5.54, 11.6, .44, 15)
    add_text_local(s, "注意：n=2/n=3 结果仅作描述性比较，不声称统计显著。", .70, 6.30, 9.0, .22, 10.5, RED, True)

    # 3. Background and objective.
    s = base_slide(prs, "Background", "背景：癌症驱动基因排序同时面对“找回已知”和“发现新候选”")
    rect(s, .70, 1.69, 3.55, 3.65, PALE_BLUE, None)
    add_text_local(s, "Recovery", 1.02, 2.01, 2.8, .30, 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "把 Validation 中已知 driver\n排进 Top-150", 1.02, 2.61, 2.8, .62, 18, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "NDCG@150 · Recall@150", 1.06, 3.53, 2.72, .26, 12, MUTED, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "优点：评价清楚\n风险：标签偏向高突变基因", 1.05, 4.16, 2.75, .56, 13, INK, align=PP_ALIGN.CENTER)
    rect(s, 4.89, 1.69, 3.55, 3.65, PALE_TEAL, None)
    add_text_local(s, "Discovery", 5.21, 2.01, 2.8, .30, 20, TEAL, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "从低频未标注基因中\n优先找到有证据候选", 5.21, 2.61, 2.8, .62, 18, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "高证据比例 · 富集倍数", 5.25, 3.53, 2.72, .26, 12, MUTED, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "优点：面向新发现\n风险：证据代理不等于真实 driver", 5.24, 4.16, 2.75, .56, 13, INK, align=PP_ALIGN.CENTER)
    rect(s, 9.08, 1.69, 3.55, 3.65, LIGHT, None)
    add_text_local(s, "核心矛盾", 9.40, 2.01, 2.8, .30, 20, ORANGE, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "一个排序很难同时把\n两类目标都做到最大", 9.40, 2.61, 2.8, .62, 18, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "需要显式表示使用者偏好", 9.44, 3.53, 2.72, .26, 12, MUTED, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "这就是“目标加权”与\n多目标强化学习的由来", 9.43, 4.16, 2.75, .56, 13, INK, align=PP_ALIGN.CENTER)
    rect(s, .72, 5.78, 11.9, .65, PALE_RED, None)
    add_text_local(s, "目标不是证明某个复杂模型最好，而是找到在偏置标签下仍可解释、稳定、可复现的取舍机制。", .98, 5.99, 11.35, .24, 14, RED, True, align=PP_ALIGN.CENTER)

    # 4. Label bias.
    s = base_slide(prs, "Why two objectives", "为什么不能只追 Recovery：已知标签天然集中在高突变频率端")
    image(s, charts["label_bias"], .58, 1.64, 7.3, 3.38)
    rect(s, 8.15, 1.70, 4.48, 3.29, LIGHT, None)
    add_text_local(s, "审计事实", 8.45, 2.02, 1.2, .22, 12, BLUE, True)
    add_rich_text(s, [
        ("Train / Validation 的突变频率百分位中位数分别为 99.8% / 98.5%。", True),
        ("静态 mutation 排序可命中 13/25 个 Validation driver。", True, RED),
        ("因此 Recovery 高分可能部分来自复现标签频率偏置，而非学习新的机制。", False),
    ], 8.42, 2.43, 3.76, 1.78, size=13.5, leading=10)
    add_text_local(s, "推论", .68, 5.42, 1.0, .22, 11, ORANGE, True)
    add_text_local(s, "需要同时报告 Discovery enrichment，并始终把 mutation 作为强 Recovery baseline；不能只用 NDCG 宣称发现能力。", .69, 5.78, 11.5, .45, 15, INK)
    source_link(s, "README.md — RL 学习回路审计与标签频率偏置", "README.md")

    # 5. Why RL was reasonable.
    s = base_slide(prs, "Initial hypothesis", "为什么最初选择序列强化学习：排序过程可能具有真实的集合历史")
    stages = [
        ("多组学 + PPI", "每个基因的\n静态证据", BLUE),
        ("当前已选集合", "覆盖患者、网络邻域、\n互补性随步骤变化", TEAL),
        ("选择下一个基因", "最大化当前增量\n而非独立打分", ORANGE),
        ("Top-150 列表", "形成有顺序的\n候选组合", NAVY),
    ]
    xs = [.67, 3.77, 6.87, 9.97]
    for i, ((head, body, color), x) in enumerate(zip(stages, xs)):
        rect(s, x, 2.06, 2.49, 2.28, LIGHT, None)
        rect(s, x, 2.06, 2.49, .11, color, color, radius=False)
        add_text_local(s, head, x+.19, 2.51, 2.10, .27, 15, color, True, align=PP_ALIGN.CENTER)
        add_text_local(s, body, x+.19, 3.10, 2.10, .66, 14, NAVY, True, align=PP_ALIGN.CENTER)
        if i < 3:
            line(s, x+2.52, 3.20, x+2.96, 3.20, GRID, 1.7)
    rect(s, .69, 5.12, 11.76, 1.02, PALE_BLUE, None)
    add_text_local(s, "合理性边界", .95, 5.36, 1.4, .22, 11, BLUE, True)
    add_text_local(s, "只有当真实 history 稳定提升排序、并且 DDQN 胜过不含 bootstrap 的 bandit 时，序列 Q-learning 才是必要的。", 2.25, 5.31, 9.75, .38, 15, NAVY, True)
    add_text_local(s, "后续实验正是围绕这两个可证伪条件展开。", 2.25, 5.79, 8.2, .24, 11, MUTED)

    # 6. Route map.
    s = base_slide(prs, "Evidence-led route", "这条路线是如何找到的：每次转向都由一个失败假设触发")
    phases = [
        ("01", "修复 RL 回路", "reward 缩放、PER、\n梯度裁剪", BLUE),
        ("02", "统一 rollout", "避免 one-pass\n偏离序列决策", TEAL),
        ("03", "scalar 权重扫描", "确认两目标\ntrade-off", ORANGE),
        ("04", "共享 MORL", "出现 conditioning /\n尺度与冲突问题", RED),
        ("05", "强化学习必要性", "即时价值模型与历史状态\n消融否定复杂性", BLUE),
        ("06", "双 head bandit", "权衡恢复，\n连续插值仍不稳", TEAL),
    ]
    for i, (num, head, body, color) in enumerate(phases):
        x = .56 + i*2.09
        rect(s, x, 2.05, 1.79, 2.55, LIGHT, None)
        rect(s, x, 2.05, 1.79, .11, color, color, radius=False)
        add_text_local(s, num, x+.14, 2.36, .42, .28, 11, color, True)
        add_text_local(s, head, x+.15, 2.78, 1.48, .47, 15, NAVY, True, align=PP_ALIGN.CENTER)
        add_text_local(s, body, x+.15, 3.52, 1.48, .57, 11.5, INK, align=PP_ALIGN.CENTER)
        if i < len(phases)-1:
            line(s, x+1.81, 3.30, x+2.03, 3.30, GRID, 1.4)
    rect(s, .70, 5.18, 11.76, .95, PALE_TEAL, None)
    add_text_local(s, "当前主线", .96, 5.45, 1.25, .22, 11, TEAL, True)
    add_text_local(s, "从“复杂模型不断补丁”切换为“最小充分模型”：先证明即时双目标价值有用，再限定成可靠的离散策略。", 2.23, 5.38, 9.72, .42, 15, NAVY, True)

    # 7. RL necessity.
    s = base_slide(prs, "Problem 1", "问题一：DDQN 的 bootstrap 没有带来稳定额外收益")
    image(s, charts["rl_necessity"], .55, 1.61, 8.1, 3.48)
    rect(s, 8.89, 1.70, 3.72, 3.30, PALE_RED, None)
    add_text_local(s, "结论", 9.17, 2.02, .9, .22, 12, RED, True)
    add_rich_text(s, [
        ("bandit 在 Recovery-heavy 区间不低于 DDQN，且方差更小。", True),
        ("DDQN 在 Discovery-heavy 区间有时更高，但不稳定、没有形成一致优势。", False),
        ("因此不能把性能归因于长期价值 bootstrap。", True, RED),
    ], 9.15, 2.43, 3.04, 1.78, size=13, leading=10)
    add_text_local(s, "解释：当前奖励主要由动作即时属性决定，未来回报可能不是主要信号。", .66, 5.52, 11.2, .31, 14, INK)
    source_link(s, "outputs/rlnecessity_20260907_152513_comparison/rl_necessity_summary.csv", "outputs/rlnecessity_20260907_152513_comparison/rl_necessity_summary.csv")

    # 8. History ablation.
    s = base_slide(prs, "Problem 2", "问题二：真实历史状态没有稳定优于 No-history / Shuffled-history")
    image(s, charts["history"], .54, 1.61, 8.12, 3.45)
    rect(s, 8.90, 1.70, 3.70, 3.28, LIGHT, None)
    add_text_local(s, "关键证据", 9.18, 2.02, 1.2, .22, 12, BLUE, True)
    add_rich_text(s, [
        ("Recovery 端：No-history 的 NDCG 0.221，高于 Full history 的 0.177。", True),
        ("Discovery 端：Full history 较高，但并未同时维持 Recovery。", False),
        ("三种模式接近或互有胜负，不支持“序列语义稳定有价值”。", True, RED),
    ], 9.15, 2.43, 3.05, 1.86, size=13, leading=9)
    add_text_local(s, "路线决策：暂停序列 Q-learning，转向 contextual bandit / ranking。", .66, 5.52, 11.2, .31, 15, NAVY, True)
    source_link(s, "outputs/historyablation_20260907_200151_comparison/history_ablation_summary.csv", "outputs/historyablation_20260907_200151_comparison/history_ablation_summary.csv")

    # 9. MORL problems.
    s = base_slide(prs, "Problem 3", "共享 MORL 的问题：条件响应、尺度与梯度修复后仍未稳定覆盖 scalar frontier")
    image(s, charts["morl_coverage"], .63, 1.64, 7.0, 3.28)
    rect(s, 7.94, 1.70, 4.65, 3.22, PALE_RED, None)
    add_text_local(s, "过程中确认的问题", 8.22, 2.01, 2.2, .24, 12, RED, True)
    add_rich_text(s, [
        ("单值Q网络容易忽略输入的目标偏好。", True),
        ("两个目标的 TD 数值尺度不同，可能出现 magnitude dominance。", False),
        ("部分 seed 的目标梯度 cosine 为负，说明共享 trunk 存在冲突。", False),
        ("vector-Q / PopArt 后仍只有 6/15 个 scalar 前沿点被覆盖。", True, RED),
    ], 8.20, 2.40, 3.87, 2.02, size=12.7, leading=8)
    add_text_local(s, "学术价值：这些负结果把问题从“继续调 reward”缩小到“模型复杂性是否必要”。", .68, 5.46, 11.2, .35, 14, INK)
    source_link(s, "outputs/morl_*_comparison/morl_vs_scalar_frontier_coverage.csv", "outputs/morl_vectorq_popart_50ep_20260904_201812_comparison/morl_vs_scalar_frontier_coverage.csv")

    # 10. New architecture.
    s = base_slide(prs, "Current method", "当前模型：可按目标偏好切换的双输出即时价值模型")
    rect(s, .67, 1.77, 3.10, 3.56, PALE_BLUE, None)
    add_text_local(s, "共享输入", .95, 2.08, 2.55, .27, 16, BLUE, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "当前状态 + 候选基因\n多组学 + 蛋白互作网络特征", .95, 2.71, 2.55, .72, 17, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "保留候选基因上下文\n但不做未来价值 bootstrap", .96, 4.08, 2.53, .52, 12.5, INK, align=PP_ALIGN.CENTER)
    line(s, 3.91, 3.45, 4.57, 3.45, BLUE, 1.8)
    rect(s, 4.77, 1.77, 3.31, 1.47, LIGHT, None)
    add_text_local(s, "找回已知基因的即时价值\nQ_rec(s,a)", 5.08, 2.13, 2.68, .54, 15, BLUE, True, align=PP_ALIGN.CENTER)
    rect(s, 4.77, 3.87, 3.31, 1.47, LIGHT, None)
    add_text_local(s, "发现新候选的即时价值\nQ_disc(s,a)", 5.08, 4.22, 2.68, .54, 15, TEAL, True, align=PP_ALIGN.CENTER)
    line(s, 8.25, 2.50, 8.89, 3.10, BLUE, 1.6)
    line(s, 8.25, 4.60, 8.89, 3.98, TEAL, 1.6)
    rect(s, 9.10, 1.77, 3.55, 3.57, PALE_TEAL, None)
    add_text_local(s, "推理时再组合", 9.40, 2.08, 2.95, .27, 16, TEAL, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "综合分数 =\n找回权重 × 找回价值\n+ 发现权重 × 发现价值", 9.43, 2.70, 2.87, 1.30, 17, NAVY, True, align=PP_ALIGN.CENTER)
    add_text_local(s, "一个模型覆盖多种取舍", 9.43, 4.49, 2.87, .28, 12.5, INK, align=PP_ALIGN.CENTER)
    rect(s, .69, 5.77, 11.91, .61, LIGHT, None)
    add_text_local(s, "训练：两个输出分别学习各自的即时奖励；不估计下一状态价值；5 个已见权重均匀调度；统一逐步生成 Top-150 排序。", .93, 5.96, 11.40, .24, 13.5, NAVY, True, align=PP_ALIGN.CENTER)

    # 11. Shared bandit positive result.
    s = base_slide(prs, "Positive result 1", "单输出共享即时价值模型：5个随机种子下已复现总体权衡方向")
    image(s, charts["shared_tradeoff"], .52, 1.58, 8.18, 3.48)
    rect(s, 8.94, 1.69, 3.66, 3.32, PALE_TEAL, None)
    add_text_local(s, "不错的结果", 9.22, 2.01, 1.4, .22, 12, TEAL, True)
    add_rich_text(s, [
        ("NDCG：0.141 → 0.094", True),
        ("Recall：0.248 → 0.176", True),
        ("新候选发现精度：0.195 → 0.260", True),
        ("富集倍数：2.10 → 2.81", True),
        ("所有 5 seed 的总体方向一致。", True, TEAL),
    ], 9.18, 2.42, 3.05, 2.00, size=13, leading=7)
    add_text_local(s, "这证明一个共享即时模型可以响应 preference，但中间权重的策略变化仍不够可控。", .66, 5.49, 11.3, .38, 14, INK)
    source_link(s, "outputs/preference_bandit_20260908_190650_5seed_extension/preference_bandit_summary.csv", "outputs/preference_bandit_20260908_190650_5seed_extension/preference_bandit_summary.csv")

    # 12. Dual-head main result.
    s = base_slide(prs, "Positive result 2", "双 head 结果：目标分离后，Recovery–Discovery 权衡更清晰")
    image(s, charts["dual_summary"], .50, 1.49, 8.22, 4.60)
    rect(s, 8.96, 1.68, 3.64, 3.95, LIGHT, None)
    add_text_local(s, "端点变化（3 seed 均值）", 9.21, 2.01, 2.95, .22, 12, NAVY, True)
    add_text_local(s, "Recovery", 9.22, 2.51, 1.25, .20, 10.5, BLUE, True)
    add_text_local(s, "NDCG   0.227 → 0.086\nRecall    0.427 → 0.160", 9.22, 2.82, 2.95, .62, 17, NAVY, True)
    line(s, 9.22, 3.64, 12.22, 3.64, GRID, .8)
    add_text_local(s, "Discovery", 9.22, 3.88, 1.35, .20, 10.5, TEAL, True)
    add_text_local(s, "发现精度  0.026 → 0.219\n富集倍数  0.283 → 2.369", 9.22, 4.19, 3.0, .62, 17, NAVY, True)
    add_text_local(s, "3/3 seed：Recovery Spearman < 0，Discovery Spearman > 0。", 9.20, 5.10, 3.0, .39, 10.5, TEAL, True)
    source_link(s, "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_summary.csv", "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_summary.csv")

    # 13. Seed-level dual-head.
    s = base_slide(prs, "Robustness", "seed 级轨迹：双 head 不是 collapse，但 Discovery 端仍有明显不稳定")
    image(s, charts["dual_seed"], .58, 1.60, 7.35, 3.53)
    rect(s, 8.18, 1.70, 4.43, 3.34, PALE_RED, None)
    add_text_local(s, "seed48 异常", 8.47, 2.01, 1.7, .22, 12, RED, True)
    add_rich_text(s, [
        ("w_disc 0.8 → 1.0 时，Recovery 反而上升：NDCG 0.160 → 0.191。", True),
        ("新候选发现反而下降：精度 0.146 → 0.098。", True, RED),
        ("纯 Discovery 的跨 seed Jaccard 仅 0.162 ± 0.170。", True, RED),
    ], 8.44, 2.43, 3.64, 1.66, size=13, leading=10)
    add_text_local(s, "解释边界：线性分数并不保证离散 Top-K 集合连续，但目标方向反转说明该端点暂不可依赖。", .68, 5.50, 11.4, .40, 14, INK)
    source_link(s, "dual_head_per_seed.csv / dual_head_across_seed_stability.csv", "outputs/dual_head_bandit_20260908_210039_comparison/dual_head_per_seed.csv")

    # 14. Smoothness.
    s = base_slide(prs, "Failure boundary", "没有通过的门槛：双 head 未改善未见 preference 的插值平滑性")
    image(s, charts["smoothness"], .54, 1.59, 7.55, 3.45)
    rect(s, 8.35, 1.70, 4.25, 3.28, PALE_RED, None)
    add_text_local(s, "双 head vs 单 head", 8.63, 2.01, 2.5, .22, 12, RED, True)
    add_rich_text(s, [
        ("相邻 Jaccard：0.736 < 0.887", True, RED),
        ("NDCG 落入相邻区间：58.3% < 66.7%", True),
        ("Recall：66.7% < 91.7%", True),
        ("Discovery：83.3% < 91.7%", True),
    ], 8.60, 2.42, 3.42, 1.65, size=13.5, leading=9)
    add_text_local(s, "因此不能宣称连续 preference 泛化；按预设门槛，不扩到 5 seed。", .67, 5.47, 11.3, .34, 15, RED, True)
    source_link(s, "interpolation_smoothness_comparison.csv", "outputs/dual_head_bandit_20260908_210039_comparison/interpolation_smoothness_comparison.csv")

    # 15. Additional visual on jumps and paired comparison.
    s = base_slide(prs, "Model comparison", "更细的证据：排序跳变集中在 Discovery 端，双 head 也未全面支配 scalar bandit")
    image(s, charts["jaccard_heatmap"], .48, 1.55, 6.18, 2.45)
    image(s, charts["paired_scalar"], 6.81, 1.55, 5.99, 2.45)
    rect(s, .67, 4.40, 12.0, 1.42, LIGHT, None)
    add_text_local(s, "独立判断", .94, 4.69, 1.2, .22, 11, BLUE, True)
    add_text_local(s, "双输出模型的价值是让两个目标分工更容易解释、权衡更明显，而不是让所有指标都更高。与每个权重单独训练的即时模型只有随机种子42/45可严格配对（n=2），不足以作显著性结论。", 2.13, 4.62, 9.95, .56, 14, NAVY, True)
    add_text_local(s, "Recovery 最佳均值 NDCG 0.235，仍低于 static mutation 0.283。", 2.13, 5.31, 8.0, .22, 11, RED, True)
    source_link(s, "dual_minus_scalar_bandit_paired.csv（paired n=2）", "outputs/dual_head_bandit_20260908_210039_comparison/dual_minus_scalar_bandit_paired.csv")

    # 16. Current decision / three profiles.
    s = base_slide(prs, "Current decision", "当前最稳妥的产品形态：三档固定策略，而不是连续滑杆")
    image(s, charts["profiles"], .56, 1.58, 6.88, 3.48)
    profiles = [
        ("Recovery-heavy", "w=(0.8,0.2)", "NDCG 0.235\nRecall 0.440", BLUE),
        ("Balanced", "w=(0.5,0.5)", "NDCG 0.222\n富集倍数 0.753", TEAL),
        ("Discovery-heavy", "w=(0.2,0.8)", "发现精度 0.156\n富集倍数 1.680", ORANGE),
    ]
    for i, (name, weight, values, color) in enumerate(profiles):
        y = 1.69 + i*1.32
        rect(s, 7.83, y, 4.74, 1.08, LIGHT, None)
        rect(s, 7.83, y, .09, 1.08, color, color, radius=False)
        add_text_local(s, name, 8.15, y+.17, 2.08, .24, 14, color, True)
        add_text_local(s, weight, 10.20, y+.18, 1.9, .22, 11, MUTED, True, align=PP_ALIGN.RIGHT)
        add_text_local(s, values, 8.15, y+.53, 3.9, .34, 12.5, NAVY, True)
    rect(s, .68, 5.55, 11.90, .59, PALE_RED, None)
    add_text_local(s, "暂不采用 w=(0,1)：纯 Discovery 端跨 seed 最不稳定；三档仍需在最终协议冻结前完成候选名单审计。", .94, 5.74, 11.38, .24, 13.5, RED, True, align=PP_ALIGN.CENTER)

    # 17. Expected outputs.
    s = base_slide(prs, "Expected outcomes", "预期成果：从“模型竞赛”转为可复现、可解释的候选生成框架")
    outcomes = [
        ("方法成果", "双目标即时价值框架\n+ 三档可解释策略", BLUE),
        ("评价成果", "已知基因找回 / 新候选发现\n双轨评价与偏置审计", TEAL),
        ("数据成果", "三套 Top-150 候选表\n跨随机种子稳定核心集合", ORANGE),
        ("学术成果", "正结果 + 负结果\n组成完整路线证据链", NAVY),
    ]
    xs = [.65, 3.73, 6.81, 9.89]
    for (head, body, color), x in zip(outcomes, xs):
        rect(s, x, 1.77, 2.52, 2.41, LIGHT, None)
        rect(s, x, 1.77, 2.52, .11, color, color, radius=False)
        add_text_local(s, head, x+.22, 2.18, 2.08, .25, 15, color, True, align=PP_ALIGN.CENTER)
        add_text_local(s, body, x+.20, 2.89, 2.12, .76, 16, NAVY, True, align=PP_ALIGN.CENTER)
    rect(s, .67, 4.72, 11.75, 1.12, PALE_TEAL, None)
    add_text_local(s, "最终可交付", .94, 4.99, 1.25, .22, 11, TEAL, True)
    add_text_local(s, "可复现实验脚本与固定协议 · 三档候选排名及证据注释 · 稳定性/帕累托图表 · 完全冻结后一次性测试集报告", 2.15, 4.91, 9.65, .46, 14.5, NAVY, True)
    add_text_local(s, "不会提前承诺“发现新 driver”；候选仍需外部数据库、文献或实验验证。", 2.15, 5.48, 8.8, .22, 10.5, RED, True)

    # 18. Next work and governance.
    s = base_slide(prs, "Next work", "后续工作：只围绕三档策略做冻结、审计和最终验证")
    tasks = [
        ("01", "冻结三档定义", "固定 (0.8,0.2) / (0.5,0.5) / (0.2,0.8)，不再连续插值。", BLUE),
        ("02", "候选稳定性审计", "每档报告跨 seed 共识基因、Jaccard、证据组成与突变频率分层。", TEAL),
        ("03", "公平基线比较", "与独立权重即时模型、图卷积网络/多层感知机、突变频率/网络度数在同一口径下比较。", ORANGE),
        ("04", "冻结后最终评估", "先锁定模型、指标和候选规则；Test 仅一次性使用且不反向决策。", RED),
    ]
    for i, (num, head, body, color) in enumerate(tasks):
        y = 1.66 + i*1.17
        rect(s, .73, y, .66, .66, color, color)
        add_text_local(s, num, .77, y+.18, .58, .18, 12, WHITE, True, align=PP_ALIGN.CENTER)
        add_text_local(s, head, 1.68, y+.04, 3.10, .24, 16, NAVY, True)
        add_text_local(s, body, 4.82, y+.03, 7.26, .43, 13.5, INK)
        if i < len(tasks)-1:
            line(s, 1.05, y+.69, 1.05, y+1.13, GRID, 1.2)
    rect(s, .72, 6.42, 11.83, .45, PALE_RED, None)
    add_text_local(s, "继续暂停：多目标强化学习修补、双深度Q网络、梯度冲突/尺度校准、训练奖励调参、基因调控网络、新癌种和额外随机种子训练。", .96, 6.54, 11.36, .20, 11.5, RED, True, align=PP_ALIGN.CENTER)

    # 19. Claims and limitations.
    s = base_slide(prs, "Claims & limitations", "可以说什么、不能说什么")
    columns = [
        (.70, "已被数据支持", TEAL, [
            "Recovery–Discovery trade-off 存在。",
            "bandit 足以复现主要权衡。",
            "双 head 能让偏好稳定改变总体方向。",
        ]),
        (4.72, "尚未被支持", RED, [
            "DDQN 有必要或优于 bandit。",
            "连续 preference 可以平滑泛化。",
            "当前模型超过所有静态基线。",
        ]),
        (8.74, "主要限制", ORANGE, [
            "最新实验仅 3 seed；scalar 配对仅 n=2。",
            "Validation 只有 25 个标签，Recall 步长为 0.04。",
            "Discovery 使用证据代理，尚非外部真实标签。",
        ]),
    ]
    for x, head, color, items in columns:
        rect(s, x, 1.72, 3.72, 4.30, LIGHT, None)
        rect(s, x, 1.72, 3.72, .12, color, color, radius=False)
        add_text_local(s, head, x+.25, 2.12, 3.15, .27, 16, color, True)
        for j, item in enumerate(items):
            bullet(s, item, x+.27, 2.76+j*.88, 3.15, .62, 13, dot=color)
    add_text_local(s, "报告原则：事实、推断和建议分开；所有模型选择均不读取 Test。", .73, 6.40, 11.6, .25, 13, NAVY, True, align=PP_ALIGN.CENTER)

    # 20. References.
    s = base_slide(prs, "Selected references", "方法背景与概念来源（点击标题可访问原文）")
    refs = [
        (
            "Roijers et al., 2013 — A Survey of Multi-Objective Sequential Decision-Making",
            "多目标序列决策需要明确偏好、scalarization 与策略覆盖集合；为本项目的 Recovery–Discovery 表述提供方法学背景。",
            "https://doi.org/10.1613/jair.3987",
            BLUE,
        ),
        (
            "Langford & Zhang, 2007 — The Epoch-Greedy Algorithm for Multi-armed Bandits with Side Information",
            "即时价值模型根据当前可观测信息选择动作并学习即时反馈，是介于监督排序与完整序列强化学习之间的简化决策模型。",
            "https://papers.nips.cc/paper_files/paper/2007/hash/4b04a686b0ad13dce35fa99fa4161c65-Abstract.html",
            TEAL,
        ),
        (
            "Hayes et al., 2022 — A Practical Guide to Multi-Objective Reinforcement Learning and Planning",
            "强调目标定义、效用形式、评价协议和覆盖集合应先于算法复杂化；与本项目的证据驱动决策门一致。",
            "https://doi.org/10.1007/s10458-022-09552-y",
            ORANGE,
        ),
    ]
    for i, (title, note, url, color) in enumerate(refs):
        y = 1.74 + i*1.48
        rect(s, .73, y, 11.86, 1.19, LIGHT, None)
        rect(s, .73, y, .09, 1.19, color, color, radius=False)
        external_link(s, title, url, 1.05, y+.20, 10.93, .29, 13.5)
        add_text_local(s, note, 1.05, y+.62, 10.90, .39, 11.5, INK)
    rect(s, .73, 6.25, 11.86, .53, PALE_BLUE, None)
    add_text_local(s, "项目数值来源均为本地 outputs/ 与 README；外部文献仅用于解释方法背景，不用于替代本项目实验结论。", 1.00, 6.41, 11.32, .23, 12, NAVY, True, align=PP_ALIGN.CENTER)


def main():
    charts = build_charts()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    title_slide(prs)
    add_slides(prs, charts)
    simplify_language(prs)
    prs.core_properties.title = "RL-GenRisk 阶段性学术汇报（2026-09-09）"
    prs.core_properties.subject = "从序列强化学习到偏好条件化双目标排序"
    prs.core_properties.author = "RL-GenRisk project"
    prs.core_properties.comments = "Generated from Train/Validation artifacts only; Test labels were not read."
    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Chart assets: {len(charts)} in {ASSET_DIR}")


if __name__ == "__main__":
    main()
